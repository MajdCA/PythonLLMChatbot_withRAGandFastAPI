import os
import json
from typing import List, Dict
import pdfplumber
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter

class DocumentLoader:
    def __init__(self, chunk_size: int = 300, chunk_overlap: int = 50):
        """
        chunk_size: Characters per chunk
        chunk_overlap: Overlap between chunks for context
        """
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
    
    def load_pdf(self, filepath: str) -> List[Dict[str, str]]:
        """Extract knowledge from PDF"""
        documents = []
        
        try:
            with pdfplumber.open(filepath) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    # Extract text
                    text = page.extract_text()
                    
                    # Extract tables (technical data)
                    tables = page.extract_tables()
                    
                    if text:
                        # Chunk the text
                        chunks = self.text_splitter.split_text(text)
                        for chunk in chunks:
                            documents.append({
                                "source": f"{os.path.basename(filepath)} (Page {page_num})",
                                "type": "pdf_text",
                                "content": chunk,
                                "page": page_num
                            })
                    
                    # Process tables separately (technical data)
                    if tables:
                        for table_idx, table in enumerate(tables):
                            table_text = self._format_table(table)
                            documents.append({
                                "source": f"{os.path.basename(filepath)} (Page {page_num}, Table {table_idx+1})",
                                "type": "pdf_table",
                                "content": table_text,
                                "page": page_num
                            })
            
            print(f"✓ Loaded {len(documents)} chunks from PDF: {filepath}")
            return documents
        
        except Exception as e:
            print(f"✗ Error loading PDF {filepath}: {str(e)}")
            return []
    
    def load_powerpoint(self, filepath: str) -> List[Dict[str, str]]:
        """Extract knowledge from PowerPoint"""
        documents = []
        
        try:
            prs = Presentation(filepath)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = ""
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                    
                    # Extract tables from shapes
                    if shape.has_table:
                        table = shape.table
                        table_text = self._format_pptx_table(table)
                        slide_text += table_text + "\n"
                
                if slide_text.strip():
                    # Chunk the slide text
                    chunks = self.text_splitter.split_text(slide_text)
                    for chunk in chunks:
                        documents.append({
                            "source": f"{os.path.basename(filepath)} (Slide {slide_num})",
                            "type": "pptx_text",
                            "content": chunk,
                            "slide": slide_num
                        })
            
            print(f"✓ Loaded {len(documents)} chunks from PowerPoint: {filepath}")
            return documents
        
        except Exception as e:
            print(f"✗ Error loading PowerPoint {filepath}: {str(e)}")
            return []
    
    def load_directory(self, directory: str) -> List[Dict[str, str]]:
        """Load all PDFs and PowerPoints from directory"""
        all_documents = []
        
        for filename in os.listdir(directory):
            filepath = os.path.join(directory, filename)
            
            if filename.lower().endswith('.pdf'):
                docs = self.load_pdf(filepath)
                all_documents.extend(docs)
            
            elif filename.lower().endswith(('.pptx', '.ppt')):
                docs = self.load_powerpoint(filepath)
                all_documents.extend(docs)
        
        return all_documents
    
    @staticmethod
    def _format_table(table: List[List[str]]) -> str:
        """Format PDF table for readability"""
        if not table:
            return ""
        
        text = "TABLE:\n"
        for row in table:
            text += " | ".join(str(cell) if cell else "" for cell in row) + "\n"
        return text
    
    @staticmethod
    def _format_pptx_table(table) -> str:
        """Format PowerPoint table for readability"""
        text = "TABLE:\n"
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            text += " | ".join(cells) + "\n"
        return text